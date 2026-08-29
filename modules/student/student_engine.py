import os
import re
import asyncio
from pathlib import Path

try:
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    docx = None

try:
    import pptx
    from pptx.util import Inches as PPTX_Inches, Pt as PPTX_Pt
    from pptx.dml.color import RGBColor as PPTX_RGBColor
except ImportError:
    pptx = None

from core.config import TEMP_DIR, logger
from modules.ai.ai_engine import ask_gemini_once
from modules.student.academic_prompts import EUROPEAN_ACADEMIC_SYSTEM_PROMPT, get_student_prompt

async def generate_student_material(user_id: int, doc_type: str, topic: str, extra_notes: str = "") -> dict:
    """
    Talabalar uchun Yevropa akademik standartidagi tayyor o'quv materialini yaratadi (.docx, .pptx yoki matn).
    """
    prompt = get_student_prompt(doc_type, topic, extra_notes)
    
    # Gemini AI yordamida ilmiy matn yaratish (1 martalik tarixsiz)
    raw_content = await ask_gemini_once(prompt)

    
    # Keraksiz AI belgilari va vizual shox-shabbalarni tozalash
    cleaned_content = re.sub(r"---|\*\*\*|===", "", raw_content).strip()
    
    file_name = f"{doc_type}_{user_id}_{topic[:15].strip().replace(' ', '_')}"
    file_name = re.sub(r"[^\w\-]", "", file_name)
    
    # PowerPoint Slayd (.pptx)
    if doc_type == "pres" and pptx:
        pptx_path = str(TEMP_DIR / f"{file_name}.pptx")
        try:
            prs = pptx.Presentation()
            
            # Titul slaydi
            blank_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = topic.title()
            subtitle.text = "Akademik Taqdimot va Ilmiy Tahlil\nTayyorlandi: Student Studio AI"
            
            # Slaydlar bo'yicha matnni bo'lish
            sections = cleaned_content.split("\n\n")
            for section in sections:
                if len(section.strip()) > 20:
                    lines = section.strip().split("\n")
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    
                    title_shape.text = lines[0].replace("#", "").strip()[:50]
                    tf = body_shape.text_frame
                    tf.text = "\n".join(lines[1:]) if len(lines) > 1 else section
                    
            prs.save(pptx_path)
            return {"type": "file", "path": pptx_path, "format": "pptx", "content": cleaned_content}
        except Exception as e:
            logger.error(f"PPTX error: {e}")
            
    # Word Hujjat (.docx) — Referat, Kurs ishi, Mustaqil ish, Maqola, Insho
    if doc_type in ["miq", "ref", "kurs", "maqola", "insho", "tezis", "test"] and docx:
        docx_path = str(TEMP_DIR / f"{file_name}.docx")
        try:
            doc = docx.Document()
            
            # Sarlavha
            title_p = doc.add_paragraph()
            title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = title_p.add_run(topic.upper() + "\n")
            run.bold = True
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(27, 54, 93)
            
            doc_type_names = {
                "miq": "MUSTAQIL ISH",
                "ref": "AKADEMIK REFERAT",
                "kurs": "KURS ISHI",
                "maqola": "ILMIY MAQOLA",
                "insho": "ESSAY / INSHO",
                "tezis": "ILMIY TEZIS",
                "test": "TEST SAVOLLARI"
            }
            subtitle_p = doc.add_paragraph()
            subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            sub_run = subtitle_p.add_run(doc_type_names.get(doc_type, "AKADEMIK ISH") + "\n\n")
            sub_run.font.size = Pt(14)
            sub_run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Matnni abzaslar bo'yicha qo'shish
            for paragraph in cleaned_content.split("\n\n"):
                p_text = paragraph.strip()
                if not p_text:
                    continue
                if p_text.startswith("#") or p_text.isupper() or len(p_text) < 60 and not p_text.endswith("."):
                    head_p = doc.add_paragraph()
                    head_run = head_p.add_run(p_text.replace("#", "").strip())
                    head_run.bold = True
                    head_run.font.size = Pt(14)
                    head_run.font.color.rgb = RGBColor(27, 54, 93)
                else:
                    body_p = doc.add_paragraph()
                    body_p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    body_run = body_p.add_run(p_text)
                    body_run.font.size = Pt(12)
                    body_run.font.name = "Times New Roman"
                    
            doc.save(docx_path)
            return {"type": "file", "path": docx_path, "format": "docx", "content": cleaned_content}
        except Exception as e:
            logger.error(f"DOCX error: {e}")
            
    return {"type": "text", "content": cleaned_content}
